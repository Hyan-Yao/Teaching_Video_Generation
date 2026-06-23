#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
make_slide.py — Generate one polished slide from arbitrary text input.

Style: navy + moss double rounded border on cream paper, bold navy title,
round-dot bullets, and a bottom 3-stage "A + B = result" diagram.

USAGE
-----
1) From a text file:
       python make_slide.py input.txt -o out.pptx

2) From stdin / a heredoc:
       python make_slide.py -o out.pptx <<'EOF'
       Title: My topic headline
       - first bullet point
       - second point with **bold** lead-in
       Diagram: Input | Process | Output
       EOF

3) From a plain string (no special markers required):
       python make_slide.py --text "Photosynthesis turns light into sugar.
       Plants capture sunlight. Chloroplasts do the work. Output is glucose." -o out.pptx

INPUT FORMAT (all parts optional — the script falls back gracefully)
--------------------------------------------------------------------
- A line starting with 'Title:'  -> slide title.
  If absent, the first sentence/line becomes the title.
- Lines starting with '-' or '*' -> bullets.
  If none are marked, the remaining sentences are auto-split into bullets.
- 'Diagram: A | B | C'           -> three captioned nodes shown as
  [A] + [B] = [C] at the bottom. If absent, the diagram is skipped and
  bullets get more vertical room.
- Inside any text, **double asterisks** mark a bold (navy) lead-in.
"""

import argparse
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------- palette -----------------------------
NAVY      = RGBColor(0x1B, 0x3A, 0x6B)
MOSS      = RGBColor(0x7A, 0x9A, 0x6B)
PAPER     = RGBColor(0xF7, 0xF6, 0xF0)
INK       = RGBColor(0x2A, 0x2A, 0x2A)
ACCENT    = RGBColor(0x2D, 0x6C, 0xDF)
ACCENT_BG = RGBColor(0xEA, 0xF0, 0xFB)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = 10.0
SLIDE_H = 7.16
FONT = "Calibri"


# ----------------------------- parsing -----------------------------
def parse_input(raw: str):
    """Turn arbitrary text into (title, [bullets], [diagram_nodes])."""
    title = None
    bullets = []
    diagram = []
    leftover = []

    for line in raw.splitlines():
        s = line.strip()
        if not s:
            continue
        low = s.lower()
        if low.startswith("title:"):
            title = s.split(":", 1)[1].strip()
        elif low.startswith("diagram:"):
            parts = s.split(":", 1)[1].split("|")
            diagram = [p.strip() for p in parts if p.strip()]
        elif s.startswith(("-", "*", "\u2022")):
            bullets.append(s.lstrip("-*\u2022 ").strip())
        else:
            leftover.append(s)

    free = " ".join(leftover).strip()
    if title is None:
        if bullets:
            title = bullets.pop(0)
        elif free:
            m = re.split(r"(?<=[.!?\u3002\uff01\uff1f])\s+", free, maxsplit=1)
            title = m[0].strip()
            free = m[1].strip() if len(m) > 1 else ""
        else:
            title = "Untitled"

    if not bullets and free:
        sentences = re.split(r"(?<=[.!?\u3002\uff01\uff1f])\s+", free)
        bullets = [s.strip() for s in sentences if s.strip()]

    return title, bullets[:7], diagram[:3]


def split_bold(text: str):
    """Split '**lead-in** rest' into runs: [(text, is_bold), ...]."""
    runs = []
    for i, chunk in enumerate(re.split(r"\*\*(.+?)\*\*", text)):
        if chunk == "":
            continue
        runs.append((chunk, i % 2 == 1))
    return runs or [(text, False)]


# --------------------------- draw helpers ---------------------------
def _no_shadow(shape):
    shape.shadow.inherit = False


def add_text(slide, x, y, w, h, runs, size, color, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None,
             font=FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = Pt(line_spacing)

    if isinstance(runs, str):
        runs = [(runs, {})]
    for txt, opts in runs:
        r = p.add_run()
        r.text = txt
        f = r.font
        f.size = Pt(opts.get("size", size))
        f.name = opts.get("font", font)
        f.bold = opts.get("bold", bold)
        f.italic = opts.get("italic", italic)
        f.color.rgb = opts.get("color", color)
    return tb


def add_round_rect(slide, x, y, w, h, fill=None, line=None, line_w=2.0,
                   radius=0.10, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    if fill is not None:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    if not shadow:
        _no_shadow(shp)
    return shp


def add_oval(slide, x, y, w, h, fill=None, line=None, line_w=2.0, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    if not shadow:
        _no_shadow(shp)
    return shp


def add_line(slide, x, y, w, color, width_pt):
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y),
                                    Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(width_pt)
    return ln


def add_arrow(slide, x, y, w, h, color):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                               Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    _no_shadow(a)
    return a


def initials(label: str) -> str:
    label = label.strip()
    if not label:
        return "?"
    if re.search(r"[\u4e00-\u9fff]", label):
        return label[:2]
    return label[0].upper()


# ----------------------------- build -------------------------------
def build(title, bullets, diagram, out_path):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = add_round_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=PAPER, radius=0.0)
    bg.line.fill.background()

    add_round_rect(slide, 0.18, 0.18, SLIDE_W - 0.36, SLIDE_H - 0.36,
                   fill=PAPER, line=NAVY, line_w=3.5, radius=0.04)
    add_round_rect(slide, 0.34, 0.34, SLIDE_W - 0.68, SLIDE_H - 0.68,
                   line=MOSS, line_w=1.75, radius=0.03)

    add_text(slide, 0.7, 0.55, 8.6, 1.15, title, size=29, color=NAVY,
             bold=True, line_spacing=33)

    add_oval(slide, 0.70, 1.74, 0.12, 0.12, fill=MOSS)
    add_line(slide, 0.84, 1.80, 4.0, MOSS, 2.5)

    bullets_bottom = 4.9 if diagram else 6.5
    top = 2.06
    n = max(len(bullets), 1)
    step = min(0.62, (bullets_bottom - top) / n)
    y = top
    for b in bullets:
        add_oval(slide, 0.82, y + 0.10, 0.10, 0.10, fill=NAVY)
        runs = [(t, {"bold": bold, "color": NAVY if bold else INK})
                for (t, bold) in split_bold(b)]
        add_text(slide, 1.06, y, 8.0, step, runs, size=15, color=INK,
                 anchor=MSO_ANCHOR.MIDDLE)
        y += step

    if diagram:
        draw_diagram(slide, diagram)

    prs.save(out_path)
    return out_path


def draw_diagram(slide, nodes):
    labels = (nodes + ["Input", "Process", "Output"])[:3]
    a, b, c = labels[0], labels[1], labels[2]
    cy = 3.95

    add_oval(slide, 1.10, cy, 1.15, 1.15, fill=WHITE, line=NAVY, line_w=2.25,
             shadow=True)
    add_text(slide, 1.10, cy, 1.15, 1.15, initials(a), size=34, color=NAVY,
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 0.55, cy + 1.22, 2.25, 0.55, a, size=13, color=NAVY,
             bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, 2.45, cy + 0.30, 0.8, 0.6, "+", size=34, color=MOSS,
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_oval(slide, 3.35, cy, 1.15, 1.15, fill=WHITE, line=MOSS, line_w=2.25,
             shadow=True)
    add_text(slide, 3.35, cy, 1.15, 1.15, initials(b), size=34, color=MOSS,
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 2.80, cy + 1.22, 2.25, 0.55, b, size=13, color=MOSS,
             bold=True, align=PP_ALIGN.CENTER)

    add_arrow(slide, 5.00, cy + 0.40, 0.6, 0.34, NAVY)

    add_round_rect(slide, 5.85, cy - 0.17, 3.3, 1.5, fill=ACCENT_BG,
                   line=ACCENT, line_w=2, radius=0.10, shadow=True)
    add_text(slide, 5.85, cy - 0.05, 3.3, 1.25,
             [(c, {"size": 22, "bold": True, "color": ACCENT})],
             size=22, color=ACCENT, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 5.85, cy + 1.40, 3.3, 0.35, "the result",
             size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ----------------------------- cli ---------------------------------
def main():
    ap = argparse.ArgumentParser(description="Generate one polished slide from text.")
    ap.add_argument("infile", nargs="?", help="input text file (or omit to read stdin)")
    ap.add_argument("--text", help="pass the text directly as a string")
    ap.add_argument("-o", "--output", default="slide.pptx", help="output .pptx path")
    args = ap.parse_args()

    if args.text is not None:
        raw = args.text
    elif args.infile:
        with open(args.infile, "r", encoding="utf-8") as f:
            raw = f.read()
    else:
        raw = sys.stdin.read()

    if not raw.strip():
        ap.error("no input text provided")

    title, bullets, diagram = parse_input(raw)
    out = build(title, bullets, diagram, args.output)
    print(f"Wrote {out}")
    print(f"  title  : {title}")
    print(f"  bullets: {len(bullets)}")
    print(f"  diagram: {diagram if diagram else '(none)'}")


if __name__ == "__main__":
    main()
